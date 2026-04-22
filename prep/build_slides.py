"""Generate the Tesla-themed PPTX for the FDE Austin case study presentation.

Design principles:
- Dark canvas (#0A0A0A), Tesla red accent (#E31937), Inter/Helvetica type.
- One insight per slide. Minimal prose. Speaker carries the narrative.
- Full speaker script lives in slide notes, not on the slide.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


BG = RGBColor(0x0A, 0x0A, 0x0A)
CARD = RGBColor(0x18, 0x18, 0x18)
CARD_HI = RGBColor(0x1E, 0x1E, 0x1E)
RED = RGBColor(0xE3, 0x19, 0x37)
RED_DIM = RGBColor(0x8E, 0x0F, 0x1F)
TEXT = RGBColor(0xF5, 0xF5, 0xF5)
MUTED = RGBColor(0x9C, 0x9C, 0x9C)
DIM = RGBColor(0x6B, 0x6B, 0x6B)
GOOD = RGBColor(0x7F, 0xB0, 0x69)
BORDER = RGBColor(0x2A, 0x2A, 0x2A)

W = Inches(13.333)
H = Inches(7.5)

FONT = "Helvetica Neue"


def _bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def _rect(slide, x, y, w, h, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.line.fill.background() if line is None else _set_line(r, line)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.shadow.inherit = False
    return r


def _set_line(shape, color, width=0.5):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def _txt(slide, x, y, w, h, text, *, size=16, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT
    return tb


def _line(slide, x1, y1, x2, y2, color=BORDER, width=0.75):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def _header(slide, eyebrow_text):
    # TESLA wordmark
    _txt(slide, Inches(0.55), Inches(0.38), Inches(3), Inches(0.35),
         "TESLA", size=11, color=RED, bold=True)
    # eyebrow / section label
    _txt(slide, Inches(1.8), Inches(0.4), Inches(8), Inches(0.3),
         eyebrow_text.upper(), size=9, color=MUTED, bold=False)
    _line(slide, Inches(0.5), Inches(0.85), Inches(12.83), Inches(0.85))


def _footer(slide, right_text):
    _line(slide, Inches(0.5), Inches(7.05), Inches(12.83), Inches(7.05))
    _txt(slide, Inches(0.55), Inches(7.15), Inches(6), Inches(0.3),
         "Athena Deployment · Austin Bldg 1", size=9, color=DIM)
    _txt(slide, Inches(7), Inches(7.15), Inches(5.83), Inches(0.3),
         right_text, size=9, color=DIM, align=PP_ALIGN.RIGHT)


def _accent_bar(slide, x, y, w, h, color=RED):
    _rect(slide, x, y, w, h, color)


def _title(slide, text, y=Inches(1.25), size=34):
    _accent_bar(slide, Inches(0.55), Inches(y.inches + 0.18), Inches(0.07), Inches(0.5), RED)
    _txt(slide, Inches(0.85), y, Inches(12), Inches(1.0), text,
         size=size, color=TEXT, bold=True)


def _subtitle(slide, text, y=Inches(2.0), size=14, color=MUTED):
    _txt(slide, Inches(0.85), y, Inches(12), Inches(0.5), text, size=size, color=color)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    return slide


# =============================================================================
# Slide builders
# =============================================================================


def slide_title(prs):
    s = add_blank(prs)
    # Full-bleed dark. Tesla wordmark large, centered-top.
    _txt(s, Inches(0.55), Inches(0.55), Inches(3), Inches(0.5),
         "TESLA", size=18, color=RED, bold=True)
    _txt(s, Inches(10), Inches(0.58), Inches(2.8), Inches(0.4),
         "CONFIDENTIAL", size=9, color=DIM, align=PP_ALIGN.RIGHT)

    # Hero title, centered-ish
    _accent_bar(s, Inches(0.55), Inches(3.0), Inches(0.1), Inches(1.6), RED)
    _txt(s, Inches(1.0), Inches(3.0), Inches(11), Inches(1.1),
         "Athena Deployment", size=56, color=TEXT, bold=True)
    _txt(s, Inches(1.0), Inches(4.1), Inches(11), Inches(0.6),
         "Austin · Building 1", size=28, color=MUTED)

    _txt(s, Inches(1.0), Inches(5.2), Inches(11), Inches(0.4),
         "FORWARD DEPLOY ENGINEER · TECHNICAL CHALLENGE", size=11, color=DIM, bold=True)

    _txt(s, Inches(0.55), Inches(6.85), Inches(10), Inches(0.3),
         "Diagnosis · Plan · Prototypes", size=10, color=MUTED)
    _txt(s, Inches(7), Inches(6.85), Inches(5.78), Inches(0.3),
         "15–20 min · Q&A 10 min", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    _notes(s, """OPEN (30 sec):
"Good [morning/afternoon]. I'm here to talk about deploying Athena at Austin Building 1 — diagnosis of what went wrong at the previous five sites, my 4-week plan, and the two features I'd prototype during the deployment.

I built this on the thesis that the prior deployments are the best available benchmark: same platform, same app teams, different outcomes. If I can understand that variance, I can design a plan that biases toward the best outcomes."

Click forward.""")
    return s


def slide_questions(prs):
    s = add_blank(prs)
    _header(s, "Framing")
    _footer(s, "2")
    _title(s, "Three questions to answer.", size=36)

    col_w = Inches(3.9)
    col_h = Inches(3.8)
    y = Inches(2.8)
    gap = Inches(0.25)
    x0 = Inches(0.85)

    questions = [
        ("01", "What went wrong\nat previous sites?"),
        ("02", "Where are the user\nand operational pain points?"),
        ("03", "How do we deploy\nAustin successfully?"),
    ]
    for i, (num, q) in enumerate(questions):
        x = Emu(x0 + Emu((col_w + gap) * i))
        _rect(s, x, y, col_w, col_h, CARD)
        _accent_bar(s, x, y, col_w, Inches(0.06), RED)
        _txt(s, Emu(x + Inches(0.4)), Emu(y + Inches(0.4)), Inches(3), Inches(0.6),
             num, size=40, color=RED, bold=True)
        _txt(s, Emu(x + Inches(0.4)), Emu(y + Inches(1.7)), Inches(3.4), Inches(1.7),
             q, size=20, color=TEXT, bold=True)

    _notes(s, """SCRIPT (30 sec):
"The challenge gave me three questions. I'll walk through what the data shows for each, then land on a 4-week Austin plan."

Don't linger. Click forward.""")
    return s


def slide_velocity(prs):
    s = add_blank(prs)
    _header(s, "What went wrong · Deployment velocity")
    _footer(s, "3")
    _title(s, "Same platform. 2× spread in time to go-live.")

    _subtitle(s, "Target: 4 weeks. Actuals across 5 prior deployments.", y=Inches(2.1))

    # Bar chart — manual rectangles, weeks scale
    # x axis origin at 3.0, target line at 7.0 (4 weeks * 1 in)
    data = [
        ("Shanghai", 3.5, GOOD),
        ("Mexico",   4.0, MUTED),
        ("Nevada",   5.0, RED_DIM),
        ("Texas",    5.0, RED_DIM),
        ("Berlin",   6.5, RED),
    ]
    chart_x = Inches(2.5)
    chart_y = Inches(2.8)
    bar_h = Inches(0.4)
    gap = Inches(0.25)
    scale_per_week = Inches(1.05)  # 1 week = 1.05"

    # target line at 4 weeks
    target_x = Emu(chart_x + Emu(scale_per_week * 4))
    _line(s, target_x, chart_y, target_x,
          Emu(chart_y + Emu((bar_h + gap) * len(data))), DIM, width=1.0)
    _txt(s, Emu(target_x - Inches(0.5)), Emu(chart_y - Inches(0.35)),
         Inches(1.8), Inches(0.25),
         "Target 4 wks", size=9, color=MUTED)

    for i, (site, weeks, color) in enumerate(data):
        y = Emu(chart_y + Emu((bar_h + gap) * i))
        # site label
        _txt(s, Inches(0.85), Emu(y + Inches(0.05)), Inches(1.55), bar_h,
             site, size=14, color=TEXT, bold=True, align=PP_ALIGN.RIGHT)
        # bar
        bar_w = Emu(scale_per_week * weeks)
        _rect(s, chart_x, y, bar_w, bar_h, color)
        # value label
        _txt(s, Emu(chart_x + bar_w + Inches(0.1)), Emu(y + Inches(0.03)),
             Inches(1.2), bar_h,
             f"{weeks:.1f} wk", size=13, color=TEXT, bold=True)

    # Key stat panels on the right
    _rect(s, Inches(10.3), Inches(2.8), Inches(2.4), Inches(1.3), CARD)
    _accent_bar(s, Inches(10.3), Inches(2.8), Inches(2.4), Inches(0.05), RED)
    _txt(s, Inches(10.45), Inches(2.9), Inches(2.2), Inches(0.3),
         "WORST", size=9, color=MUTED, bold=True)
    _txt(s, Inches(10.45), Inches(3.2), Inches(2.2), Inches(0.5),
         "Berlin +62.5%", size=20, color=RED, bold=True)
    _txt(s, Inches(10.45), Inches(3.65), Inches(2.2), Inches(0.4),
         "6.5 wks vs 4 target", size=11, color=MUTED)

    _rect(s, Inches(10.3), Inches(4.25), Inches(2.4), Inches(1.3), CARD)
    _accent_bar(s, Inches(10.3), Inches(4.25), Inches(2.4), Inches(0.05), GOOD)
    _txt(s, Inches(10.45), Inches(4.35), Inches(2.2), Inches(0.3),
         "BEST", size=9, color=MUTED, bold=True)
    _txt(s, Inches(10.45), Inches(4.65), Inches(2.2), Inches(0.5),
         "Shanghai −12.5%", size=20, color=GOOD, bold=True)
    _txt(s, Inches(10.45), Inches(5.1), Inches(2.2), Inches(0.4),
         "3.5 wks ahead of plan", size=11, color=MUTED)

    # Bottom takeaway
    _rect(s, Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.7), CARD_HI)
    _txt(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.5),
         "Same platform. Same 14 engineers. The deployment process is the variable.",
         size=15, color=TEXT, bold=True)

    _notes(s, """SCRIPT (2 min):
"Start here — five prior Athena deployments, target was 4 weeks each time. Actuals range from 3.5 at Shanghai to 6.5 at Berlin. Two out of five hit the target.

What matters is the SPREAD. The platform is the same. The engineers are the same. But outcomes vary by 2x. That tells me there's a deployment PROCESS variable — sequencing, stakeholder engagement, integration strategy — that I can control.

So the question stops being 'is Athena ready' and becomes 'what is Shanghai doing right, and what broke at Berlin?' The rest of my presentation answers that."

Don't call Berlin a failure. "Over target" is the right phrase. Several panelists may have run Berlin.""")
    return s


def slide_berlin(prs):
    s = add_blank(prs)
    _header(s, "What went wrong · Where the time went")
    _footer(s, "4")
    _title(s, "Berlin: 19 delay-days. $192K disruption.")

    _subtitle(s, "Integration and compliance — not app bugs — drove the overrun.",
              y=Inches(2.1))

    # 3 blocker cards
    blockers = [
        ("GDPR audit", "4 days", "Compliance wasn't engaged early; change-tracking audit logs flagged."),
        ("SAP integration", "5 days", "Budget codes didn't reconcile with Athena Area/System IDs."),
        ("Primavera sync", "5 days", "Legacy schedule milestones incompatible with Athena milestones."),
    ]
    x0 = Inches(0.85)
    card_w = Inches(3.85)
    card_h = Inches(2.6)
    gap = Inches(0.3)
    y = Inches(2.9)

    for i, (name, days, desc) in enumerate(blockers):
        x = Emu(x0 + Emu((card_w + gap) * i))
        _rect(s, x, y, card_w, card_h, CARD)
        _accent_bar(s, x, y, Inches(0.06), card_h, RED)
        _txt(s, Emu(x + Inches(0.35)), Emu(y + Inches(0.3)), Inches(3.3), Inches(0.4),
             name.upper(), size=11, color=MUTED, bold=True)
        _txt(s, Emu(x + Inches(0.35)), Emu(y + Inches(0.75)), Inches(3.3), Inches(0.9),
             days, size=40, color=RED, bold=True)
        _txt(s, Emu(x + Inches(0.35)), Emu(y + Inches(1.65)), Inches(3.3), Inches(0.85),
             desc, size=11, color=TEXT)

    _rect(s, Inches(0.85), Inches(5.9), Inches(11.5), Inches(0.8), CARD_HI)
    _txt(s, Inches(1.0), Inches(6.02), Inches(11.3), Inches(0.6),
         "Takeaway: integration + compliance have the longest tails. Front-load both in week 1.",
         size=15, color=TEXT, bold=True)

    _notes(s, """SCRIPT (2 min):
"Drilling into Berlin — the worst site — 19 of the total delay-days happened here, almost $192K in disruption cost. Eight distinct blocking events.

The top three blockers I'd call out: GDPR audit flagged change-tracking logs four days in; SAP budget codes didn't match Athena's Area/System IDs — five-day block; Primavera milestone sync — five-day block, and Nevada hit the same issue, so this is a pattern.

The takeaway for Austin: integration and compliance drove the longest tails, NOT Athena app bugs. The mitigation is to front-load both of those in week 1 — SAP reconciliation on Tuesday, Primavera sync on Wednesday, Compliance kickoff Monday with F. Mueller. I'll show the week 1 plan shortly."

If asked: "how do you know it wasn't the app" — because the issue records show integration failures (SAP codes, Primavera sync) and regulatory pushback (GDPR, GC audit), not defects in Athena core features.""")
    return s


def slide_themes(prs):
    s = add_blank(prs)
    _header(s, "User pain · Cross-site themes")
    _footer(s, "5")
    _title(s, "Pain is platform-wide, not site-specific.")

    _subtitle(s, "38 feedback items across 5 sites. Every top theme appears at ≥4 of them.",
              y=Inches(2.1))

    themes = [
        ("01", "Change-management clicks & workflow", 8, 5),
        ("02", "Tekton ↔ Athena sync (esp. offline)",  8, 4),
        ("03", "Milestones & scheduling",              7, 3),
        ("04", "Reporting, variance, rollups",          7, 2),
    ]
    y = Inches(2.95)
    row_h = Inches(0.85)
    for i, (num, name, mentions, hi) in enumerate(themes):
        cy = Emu(y + Emu(row_h * i))
        _rect(s, Inches(0.85), cy, Inches(11.5), Inches(0.75), CARD if i % 2 == 0 else CARD_HI)
        _accent_bar(s, Inches(0.85), cy, Inches(0.06), Inches(0.75), RED)
        _txt(s, Inches(1.1), Emu(cy + Inches(0.2)), Inches(0.7), Inches(0.5),
             num, size=18, color=RED, bold=True)
        _txt(s, Inches(1.9), Emu(cy + Inches(0.2)), Inches(7), Inches(0.5),
             name, size=17, color=TEXT, bold=True)
        _txt(s, Inches(9.0), Emu(cy + Inches(0.2)), Inches(1.6), Inches(0.5),
             f"{mentions}  mentions", size=14, color=TEXT, align=PP_ALIGN.RIGHT)
        _txt(s, Inches(10.7), Emu(cy + Inches(0.22)), Inches(1.6), Inches(0.5),
             f"{hi} high-pri", size=12, color=MUTED, align=PP_ALIGN.RIGHT)

    _rect(s, Inches(0.85), Inches(6.55), Inches(11.5), Inches(0.35), CARD_HI)
    _txt(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.3),
         "22 of 38 feedback items have no status recorded — triage itself is broken.",
         size=12, color=RED, bold=True)

    _notes(s, """SCRIPT (2.5 min):
"I tagged every feedback item across all five sites into themes. Four bubble up.

The #1 theme is change-management — the clicks, the workflow, the approval chain. Quote from the data: 'change management workflow is unclear — takes 10-plus clicks to adjust scope hours.' That appears at every site.

#2 is Tekton-Athena sync — especially offline mode in the field. Mexico, Texas, Berlin all had blockers here.

#3 and #4 tie — milestones / scheduling and reporting / variance rollups.

The important line: every one of these shows up at 4 of 5 sites or more. That means this is not a Berlin problem or a Mexico problem. These are platform gaps. The FDE role isn't to fix the platform — it's to deploy around them and feed the backlog with concrete prototypes, which I'll get to.

Red text at the bottom: 22 of 38 feedback items have no status recorded. Most of the 'unresolved high-priority' items were hiding in the untriaged pile. That's a process failure — nobody owns the triage. The fix for Austin: FDE owns weekly triage of the feedback log, mandatory status field."

Common follow-up: "How did you tag themes?" — keyword matching against the pain_point field plus manual spot-check. Not perfect, but the top four are clear signal.""")
    return s


def slide_adoption(prs):
    s = add_blank(prs)
    _header(s, "Adoption friction")
    _footer(s, "6")
    _title(s, "Shanghai & Mexico hit zero blocked users by week 3.")

    _subtitle(s, "Berlin was still blocking 14% of trained users at week 3.",
              y=Inches(2.1))

    # Two side-by-side callouts
    _rect(s, Inches(0.85), Inches(2.9), Inches(5.75), Inches(3.5), CARD)
    _accent_bar(s, Inches(0.85), Inches(2.9), Inches(5.75), Inches(0.06), GOOD)
    _txt(s, Inches(1.1), Inches(3.05), Inches(5.4), Inches(0.4),
         "WHAT SHANGHAI / MEXICO DID", size=10, color=MUTED, bold=True)
    _txt(s, Inches(1.1), Inches(3.4), Inches(5.4), Inches(0.6),
         "Fast adoption ramp", size=22, color=TEXT, bold=True)
    bullets_left = [
        "Shipped a real feature in week 2 (Shanghai: change-order baseline · 2-day build)",
        "Small training cohorts (40–50 users)",
        "Week 1 internal advocate (prior-site PM)",
        "Tactical wins during deploy built trust",
    ]
    for i, b in enumerate(bullets_left):
        _txt(s, Inches(1.1), Emu(Inches(4.2) + Emu(Inches(0.45) * i)),
             Inches(5.4), Inches(0.4),
             f"·  {b}", size=11, color=TEXT)

    _rect(s, Inches(6.8), Inches(2.9), Inches(5.55), Inches(3.5), CARD)
    _accent_bar(s, Inches(6.8), Inches(2.9), Inches(5.55), Inches(0.06), RED)
    _txt(s, Inches(7.05), Inches(3.05), Inches(5.2), Inches(0.4),
         "WHAT BERLIN DID", size=10, color=MUTED, bold=True)
    _txt(s, Inches(7.05), Inches(3.4), Inches(5.2), Inches(0.6),
         "Sticky adoption friction", size=22, color=RED, bold=True)
    bullets_right = [
        "22 users blocked week 1 (vs 10 at Shanghai)",
        "Compliance engaged late (week 3)",
        "Larger training cohorts, less hands-on",
        "Integration issues compounded user blocks",
    ]
    for i, b in enumerate(bullets_right):
        _txt(s, Inches(7.05), Emu(Inches(4.2) + Emu(Inches(0.45) * i)),
             Inches(5.2), Inches(0.4),
             f"·  {b}", size=11, color=TEXT)

    _rect(s, Inches(0.85), Inches(6.55), Inches(11.5), Inches(0.35), CARD_HI)
    _txt(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.3),
         "Austin plan borrows from Shanghai: internal champion, small cohorts, week-3 prototype ship.",
         size=12, color=TEXT, bold=True)

    _notes(s, """SCRIPT (1 min):
"Adoption metric I care about: blocked users over time. Users who can't do their job because of a config, permission, or training gap.

Shanghai and Mexico drop to zero by week 3. Berlin is still at 14% blocked at week 3 and doesn't reach zero until week 6.

Why the difference? Three things I'd copy from Shanghai:
  - Ship a real, tactical feature during deployment — builds trust fast.
  - Small training cohorts (40 to 50 users max) with hands-on practice.
  - Named internal champion in week 1 — someone who's used Athena before.

Austin has the ideal candidate: P. Kumar, senior PM with deep Athena knowledge from a prior site. I'll pair him with skeptical teams in week 1."

If asked: "why did Berlin have more blocks?" — the data doesn't tell me precisely, but the event log shows SAP + GDPR blockers in week 2 that cascaded into permissions issues. My inference: fix the integrations first and the adoption friction drops naturally.""")
    return s


def slide_austin_brief(prs):
    s = add_blank(prs)
    _header(s, "Austin · Situation brief")
    _footer(s, "7")
    _title(s, "What we're deploying into.")

    _subtitle(s, "Building 1 · from the estimate and technical survey.", y=Inches(2.1))

    stats = [
        ("36", "scopes", "across 4 Areas/Systems · ~2,800 tasks"),
        ("85", "Tekton users", "420 task-hour entries per day"),
        ("7", "external parties", "5 contractors + 2 design partners · RBAC day 1"),
        ("10%", "API-fluent", "25% PM-software · 60% Primavera · 40% SAP"),
    ]
    x0 = Inches(0.85)
    y = Inches(2.9)
    card_w = Inches(2.85)
    card_h = Inches(2.6)
    gap = Inches(0.23)
    for i, (num, label, sub) in enumerate(stats):
        x = Emu(x0 + Emu((card_w + gap) * i))
        _rect(s, x, y, card_w, card_h, CARD)
        _accent_bar(s, x, y, card_w, Inches(0.06), RED)
        _txt(s, Emu(x + Inches(0.3)), Emu(y + Inches(0.4)), Inches(2.55), Inches(1.3),
             num, size=60, color=TEXT, bold=True)
        _txt(s, Emu(x + Inches(0.3)), Emu(y + Inches(1.55)), Inches(2.55), Inches(0.5),
             label, size=14, color=RED, bold=True)
        _txt(s, Emu(x + Inches(0.3)), Emu(y + Inches(1.95)), Inches(2.55), Inches(0.7),
             sub, size=10, color=MUTED)

    _rect(s, Inches(0.85), Inches(5.8), Inches(11.5), Inches(0.95), CARD_HI)
    _accent_bar(s, Inches(0.85), Inches(5.8), Inches(0.06), Inches(0.95), RED)
    _txt(s, Inches(1.1), Inches(5.92), Inches(11.2), Inches(0.35),
         "BIGGEST RISK", size=10, color=MUTED, bold=True)
    _txt(s, Inches(1.1), Inches(6.22), Inches(11.2), Inches(0.5),
         "Zones C/D: 2–3 hrs offline per day. Tekton offline mode hasn't been hardened above 50 tasks.",
         size=14, color=TEXT, bold=True)

    _notes(s, """SCRIPT (1.5 min):
"Switching from diagnosis to what I'd walk into at Austin.

Four numbers from the estimate and technical survey:
- 36 scopes across 4 Areas/Systems in Building 1, roughly 2,800 tasks. That's the breakout I'd load into Athena day one.
- 85 Tekton field users, 420 task-hour entries per day. Non-trivial throughput.
- 7 external parties needing scope access: five contractors and two design partners. Every one of them needs RBAC configured before they can log in.
- 10% API fluency. 25% PM software. That's the training gap — hands-on, not docs.

Biggest Day-1 risk bar none: zones C and D have 2 to 3 hours of offline time per day per the tech survey. Tekton offline mode has been prototyped for 50 tasks; we need it to work for hundreds. That's one of my two prototype picks — I'll get there."

If asked: "where does the 2,800 task number come from?" — technical survey, page of notes from R. Patel on Oct 4.""")
    return s


def slide_plan(prs):
    s = add_blank(prs)
    _header(s, "4-week plan")
    _footer(s, "8")
    _title(s, "Deploy Austin in 4 weeks. Exit with a playbook.")

    weeks = [
        ("WK 1", "Infra · integrations · data load",
         ["Primavera + SAP wired in staging",
          "Austin estimate loaded (defects fixed)",
          "RBAC for 7 external parties",
          "Compliance kickoff (F. Mueller)"]),
        ("WK 2", "Pilot team + Tekton field trial",
         ["1 Area/System live (Electrical)",
          "Tekton trial in zone C",
          "P. Kumar pair-instructing",
          "Defect retro Friday"]),
        ("WK 3", "Full rollout + prototype 1",
         ["3 Areas/Systems go live",
          "Ship 1-click change approval",
          "45-min training cohorts",
          "Weekly report to R. Lopez"]),
        ("WK 4", "Audit · prototype 2 · handoff",
         ["Compliance audit Wednesday",
          "Ship hardened Tekton offline",
          "Playbook to next FDE",
          "Exit review"]),
    ]
    x0 = Inches(0.85)
    y = Inches(2.25)
    card_w = Inches(2.85)
    card_h = Inches(4.25)
    gap = Inches(0.23)
    for i, (wk, focus, exits) in enumerate(weeks):
        x = Emu(x0 + Emu((card_w + gap) * i))
        _rect(s, x, y, card_w, card_h, CARD)
        _accent_bar(s, x, y, card_w, Inches(0.08), RED)
        _txt(s, Emu(x + Inches(0.3)), Emu(y + Inches(0.3)), Inches(2.55), Inches(0.45),
             wk, size=13, color=RED, bold=True)
        _txt(s, Emu(x + Inches(0.3)), Emu(y + Inches(0.85)), Inches(2.55), Inches(1.15),
             focus, size=15, color=TEXT, bold=True)
        for j, e in enumerate(exits):
            _txt(s, Emu(x + Inches(0.3)),
                 Emu(y + Inches(2.2) + Emu(Inches(0.42) * j)),
                 Inches(2.55), Inches(0.5),
                 f"·  {e}", size=10, color=TEXT)

    _txt(s, Inches(0.85), Inches(6.7), Inches(11.5), Inches(0.3),
         "Stakeholder constraints built-in: K. Abrams traveling Wk2 · D. Nakamura oncall Wk3 · Audit Wk4.",
         size=11, color=MUTED)

    _notes(s, """SCRIPT (3.5 min — the core slide):
"This is the plan. Four weeks, one exit criteria per week.

Week 1 is integration-heavy because Berlin told me that's where the longest tails are. Primavera and SAP wired and tested in staging by Friday. Estimate loaded — and I mean CLEAN, I'll show the defects I'd fix in a moment. All seven external parties provisioned. Compliance kickoff Monday — not week 4 — F. Mueller gets briefed early.

Week 2 is the pilot. One Area/System, probably Electrical because it has the most scopes so we learn fast. P. Kumar pair-instructs with the pilot team. Tekton field trial in zone C — the worst offline zone — in PRODUCTION, not staging. Friday defect retro. If anything's still on fire going into the weekend, I'd rather know.

Week 3 is full rollout. Three remaining Areas/Systems in parallel. 45-minute hands-on training cohorts, 40 to 50 users each, Shanghai-style. This is when I ship the first prototype — 1-click change approval for scopes under $5K — which is a back-port of a prototype that already worked at Texas. And we hit the first weekly report to R. Lopez on Thursday.

Week 4 is audit week. Audit dry-run Monday with F. Mueller so we catch anything before the real audit Wednesday. Second prototype ships: hardened Tekton offline mode. Handoff playbook Friday.

One thing I want to flag: K. Abrams owns milestone approvals and he's traveling Wednesday through Friday of sprint 2 — I'd pre-approve any sprint 2 milestones Monday and delegate to S. Chen. D. Nakamura is oncall in sprint 3, so the Thursday summary is hard-locked, no mid-week asks. And the audit in sprint 4 is known — I plan around it, not into it."

Have Plan B ready: "if only 2 weeks, drop full training, go pair-instruct with Kumar, skip pilot, ship only prototype 1, flag trade-off to R. Lopez upfront."
""")
    return s


def slide_prototypes(prs):
    s = add_blank(prs)
    _header(s, "Prototype picks")
    _footer(s, "9")
    _title(s, "Two prototypes in weeks 3 and 4.")

    _subtitle(s, "Both back-ported from features other sites already validated. Low risk, high leverage.",
              y=Inches(2.1))

    protos = [
        (
            "01",
            "1-click change approval  <$5K",
            "Back-port from Texas prototype (already working for cost changes).",
            [
                "Addresses #1 pain theme (change-mgmt clicks)",
                "Clips ~80% of approval volume by value",
                "2–3 day build · proven pattern",
            ],
        ),
        (
            "02",
            "Tekton offline + batch auto-queue",
            "Harden the Texas prototype from 50 → 500+ tasks.",
            [
                "Critical for zones C/D at Austin (2–3 hrs offline/day)",
                "Berlin lost 8 hrs of entries — we will too without it",
                "4–5 day build · mobile + backend sync",
            ],
        ),
    ]
    x0 = Inches(0.85)
    y = Inches(2.85)
    card_w = Inches(5.8)
    card_h = Inches(3.7)
    gap = Inches(0.3)
    for i, (num, title, sub, bullets) in enumerate(protos):
        x = Emu(x0 + Emu((card_w + gap) * i))
        _rect(s, x, y, card_w, card_h, CARD)
        _accent_bar(s, x, y, Inches(0.08), card_h, RED)
        _txt(s, Emu(x + Inches(0.4)), Emu(y + Inches(0.3)), Inches(1), Inches(0.5),
             num, size=11, color=RED, bold=True)
        _txt(s, Emu(x + Inches(0.4)), Emu(y + Inches(0.65)), Inches(5.3), Inches(0.7),
             title, size=20, color=TEXT, bold=True)
        _txt(s, Emu(x + Inches(0.4)), Emu(y + Inches(1.5)), Inches(5.3), Inches(0.55),
             sub, size=12, color=MUTED)
        for j, b in enumerate(bullets):
            _txt(s, Emu(x + Inches(0.4)),
                 Emu(y + Inches(2.2) + Emu(Inches(0.42) * j)),
                 Inches(5.3), Inches(0.4),
                 f"·  {b}", size=11, color=TEXT)

    _rect(s, Inches(0.85), Inches(6.75), Inches(11.5), Inches(0.3), CARD_HI)
    _txt(s, Inches(1.0), Inches(6.78), Inches(11.3), Inches(0.25),
         "Explicitly deferred: photo upload (8 wks), cost forecasting (needs data scientist), RBAC overhaul (4 wks).",
         size=10, color=MUTED)

    _notes(s, """SCRIPT (1.5 min):
"Two features I'd prototype during the deployment, not separately.

One: 1-click change approval for scopes under 5K dollars. Texas already prototyped this and it works. We'd forward-port. This directly hits my #1 pain theme — change-management clicks — and clips the vast majority of approval volume without touching the bigger-change flow. 2 to 3 day build.

Two: Tekton offline mode with batch auto-queue. The Texas prototype works for 50 tasks. We need it for 500-plus because Austin has 85 field users averaging 420 entries a day in zones with known offline exposure. Berlin lost 8 hours of field data to this exact issue. 4 to 5 day build.

What I'm explicitly NOT doing: photo upload integration — 8-week project, can't fit it. Cost forecasting dashboard — needs a data scientist and historical burn-rate data we don't have. Full RBAC overhaul — 4-week project, touches auth. All three are legitimate but they're week 8-plus work. I'd add them to the post-go-live backlog with explicit next-FDE ownership."

Common follow-up: "could you ship these faster?" — yes, I could ship only one and compress. Trade-off is we hit zero-offline in zone C on day one without the batch queue hardening, which means lost data. Not worth it unless we're genuinely two weeks.""")
    return s


def slide_day1(prs):
    s = add_blank(prs)
    _header(s, "Austin · Day-1 data quality")
    _footer(s, "10")
    _title(s, "Fix 41 defects before the Athena import.")

    _subtitle(s, "Found in the Austin Estimate — the file that defines Area/System/Scope structure.",
              y=Inches(2.1))

    # Big stat
    _rect(s, Inches(0.85), Inches(2.9), Inches(4.5), Inches(2.6), CARD)
    _accent_bar(s, Inches(0.85), Inches(2.9), Inches(4.5), Inches(0.06), RED)
    _txt(s, Inches(1.05), Inches(3.0), Inches(4.2), Inches(0.3),
         "DEFECTS FOUND", size=10, color=MUTED, bold=True)
    _txt(s, Inches(1.05), Inches(3.3), Inches(4.2), Inches(1.2),
         "41", size=88, color=RED, bold=True)
    _txt(s, Inches(1.05), Inches(4.55), Inches(4.2), Inches(0.5),
         "cells across 6 defect types", size=14, color=TEXT)
    _txt(s, Inches(1.05), Inches(4.95), Inches(4.2), Inches(0.5),
         "in 36 scope rows", size=11, color=MUTED)

    # Right: defect list
    defects = [
        ("Project", "Missing", 1),
        ("Area/System", "Missing (scope would orphan)", 1),
        ("Labor Cost", "Mixed $ / commas / k suffixes", 13),
        ("Material Cost", "Mixed formats", 6),
        ("Rental Cost", "Mixed formats", 2),
        ("Material Qty", "Mixed unit formats", 18),
    ]
    x_list = Inches(5.8)
    y_list = Inches(2.9)
    _rect(s, x_list, y_list, Inches(6.55), Inches(3.85), CARD)
    for i, (field, defect, rows) in enumerate(defects):
        ry = Emu(y_list + Inches(0.3) + Emu(Inches(0.52) * i))
        _txt(s, Emu(x_list + Inches(0.3)), ry, Inches(1.8), Inches(0.4),
             field, size=12, color=RED, bold=True)
        _txt(s, Emu(x_list + Inches(2.1)), ry, Inches(3.3), Inches(0.4),
             defect, size=12, color=TEXT)
        _txt(s, Emu(x_list + Inches(5.6)), ry, Inches(0.8), Inches(0.4),
             str(rows), size=14, color=TEXT, bold=True, align=PP_ALIGN.RIGHT)

    _rect(s, Inches(0.85), Inches(6.85), Inches(11.5), Inches(0.2), CARD_HI)
    _txt(s, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.25),
         "30-min session with Estimating team · before kickoff",
         size=10, color=MUTED)

    _notes(s, """SCRIPT (1 min):
"The challenge hint says 'the Athena project breakout is defined by a provided estimate.' So this file IS the structure.

I ran a defect scan. 41 cells, 6 defect types, in only 36 scope rows. One row is missing its Project field. One row has no Area/System — that scope would orphan on import. And labor/material/rental costs are a mess of dollar signs, commas, and 'k' or 'M' suffixes — 39 cells across those columns.

If we load this unchanged, one of two things happens: the numeric parser crashes, or worse, it silently miscategorizes.

The fix: 30 minutes with the Estimating team before kickoff. Lock the format, re-export, validate. Then we're clean."

If asked: "is this blocking?" — yes for the malformed cost columns; a defensive parser might save us but you wouldn't want to discover a numeric format bug during a live load.""")
    return s


def slide_stakeholders(prs):
    s = add_blank(prs)
    _header(s, "Stakeholders · engagement model")
    _footer(s, "11")
    _title(s, "People. Cadence. Champions.")

    _subtitle(s, "Named owners, explicit cadences, no ambiguity.", y=Inches(2.1))

    rows = [
        ("DAILY", "9am standup", "All project teams + FDE. Mandatory per contact sheet."),
        ("WEEKLY", "Thursday report → R. Lopez", "Project Controls Lead. She asked for it."),
        ("WEEK 1", "Compliance kickoff → F. Mueller", "NOT week 4. Berlin learned this the hard way."),
        ("WEEK 1+", "Pair with P. Kumar (internal champion)", "Prior-site Athena veteran; good advocate."),
        ("WEEKLY", "External partners sync", "5 contractors + 2 design partners. Their SLA is 2 wks."),
    ]
    x0 = Inches(0.85)
    y0 = Inches(2.85)
    row_h = Inches(0.75)
    for i, (when, who, note) in enumerate(rows):
        y = Emu(y0 + Emu(row_h * i))
        _rect(s, x0, y, Inches(11.5), Inches(0.68), CARD if i % 2 == 0 else CARD_HI)
        _accent_bar(s, x0, y, Inches(0.06), Inches(0.68), RED)
        _txt(s, Inches(1.1), Emu(y + Inches(0.2)), Inches(1.3), Inches(0.4),
             when, size=10, color=RED, bold=True)
        _txt(s, Inches(2.5), Emu(y + Inches(0.17)), Inches(5.5), Inches(0.4),
             who, size=15, color=TEXT, bold=True)
        _txt(s, Inches(7.8), Emu(y + Inches(0.2)), Inches(4.5), Inches(0.4),
             note, size=11, color=MUTED)

    _notes(s, """SCRIPT (1 min):
"Engagement model — four touchpoints I'd lock in day one.

Daily 9am standup — it's already mandatory per the contact sheet, I'd just make sure FDE shows up every day. That's the blast radius for status issues.

Weekly Thursday report to R. Lopez, Project Controls Lead — she specifically asked for it. She's also the one who'll push on my composite risk score, so that weekly cadence is how trust gets built.

F. Mueller, Compliance — brought in WEEK 1, not week 4. Berlin hit the GDPR wall because Compliance wasn't briefed on change-tracking requirements until it was too late.

P. Kumar is my internal champion. Senior PM, deep Athena knowledge from a prior site, happy to help other teams. He's paired with skeptical teams in week 1.

External partners — weekly sync. Their change-request SLA is 2 weeks. I can't move faster than them so I plan for that latency."

Common question: "what if R. Lopez pushes back on the plan?" — invite her to the Monday kickoff; the plan has her check-ins baked in; treat her as co-owner.""")
    return s


def slide_recap(prs):
    s = add_blank(prs)
    _header(s, "Recap")
    _footer(s, "12")
    _title(s, "What I'm committing to.", size=36)

    items = [
        ("Deploy in 4 weeks", "On or under target. Shanghai is the benchmark."),
        ("Two prototypes", "1-click change approval · hardened Tekton offline."),
        ("Exit with a playbook", "So the next FDE starts ahead of where I started."),
    ]
    y0 = Inches(2.85)
    for i, (head, sub) in enumerate(items):
        y = Emu(y0 + Emu(Inches(1.3) * i))
        _rect(s, Inches(0.85), y, Inches(11.5), Inches(1.15), CARD)
        _accent_bar(s, Inches(0.85), y, Inches(0.08), Inches(1.15), RED)
        _txt(s, Inches(1.1), Emu(y + Inches(0.2)), Inches(1), Inches(0.5),
             f"0{i+1}", size=11, color=RED, bold=True)
        _txt(s, Inches(2.0), Emu(y + Inches(0.17)), Inches(10), Inches(0.55),
             head, size=24, color=TEXT, bold=True)
        _txt(s, Inches(2.0), Emu(y + Inches(0.65)), Inches(10), Inches(0.45),
             sub, size=13, color=MUTED)

    _txt(s, Inches(0.85), Inches(6.7), Inches(11.5), Inches(0.4),
         "Happy to dig into any of this.", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    _notes(s, """SCRIPT (30 sec):
"To recap — three commitments.

One: deploy Austin in four weeks. Shanghai's 3.5 is the benchmark, not the ceiling. If I beat Shanghai, great. If I hit four, that's the target.

Two: ship two prototypes during the deployment — 1-click change approval and hardened Tekton offline. Both validated by prior sites. Low risk.

Three: I exit with a playbook. So the next FDE inherits a process, not a recovery.

Happy to dig into any of this."

Then STOP. Let them ask.""")
    return s


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_questions(prs)
    slide_velocity(prs)
    slide_berlin(prs)
    slide_themes(prs)
    slide_adoption(prs)
    slide_austin_brief(prs)
    slide_plan(prs)
    slide_prototypes(prs)
    slide_day1(prs)
    slide_stakeholders(prs)
    slide_recap(prs)

    out = Path(__file__).parent / "Athena_Austin_FDE.pptx"
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
